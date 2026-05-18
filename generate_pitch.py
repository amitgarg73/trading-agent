"""
Generates Trading_Agent_Pitch_Deck.pptx — VC-ready pitch deck.
Covers value prop, market opportunity, MOAT, traction, and business model
without exposing technical mechanics.

Run: python3 generate_pitch.py
"""
import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

# ── Color palette ─────────────────────────────────────────────────────────────
NAVY       = RGBColor(0x1A, 0x3A, 0x6A)
ORANGE     = RGBColor(0xF4, 0x7B, 0x20)
WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RGBColor(0xF2, 0xF4, 0xF7)
MID_GRAY   = RGBColor(0x88, 0x88, 0x99)
GREEN      = RGBColor(0x1A, 0x7A, 0x3A)
DARK_GRAY  = RGBColor(0x33, 0x33, 0x44)

SLIDE_W = Inches(13.33)
SLIDE_H = Inches(7.5)

prs = Presentation()
prs.slide_width  = SLIDE_W
prs.slide_height = SLIDE_H

blank_layout = prs.slide_layouts[6]  # completely blank


# ── Helpers ───────────────────────────────────────────────────────────────────

def add_rect(slide, l, t, w, h, fill_rgb, line_rgb=None, line_width_pt=0):
    shape = slide.shapes.add_shape(1, Inches(l), Inches(t), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_rgb
    if line_rgb:
        shape.line.color.rgb = line_rgb
        shape.line.width = Pt(line_width_pt)
    else:
        shape.line.fill.background()
    return shape


def add_text(slide, text, l, t, w, h,
             font_size=18, bold=False, italic=False,
             color=WHITE, align=PP_ALIGN.LEFT,
             word_wrap=True, font_name="Calibri"):
    txBox = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = txBox.text_frame
    tf.word_wrap = word_wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = font_name
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return txBox


def add_multi_para(slide, lines, l, t, w, h,
                   font_size=14, color=WHITE, bullet_char="•",
                   line_spacing_pt=4, font_name="Calibri"):
    """Add multiple bullet lines to one text box."""
    txBox = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = txBox.text_frame
    tf.word_wrap = True
    first = True
    for line in lines:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        if bullet_char and not line.startswith("__"):
            p.text = f"{bullet_char}  {line}"
        else:
            p.text = line.lstrip("_")
        for run in p.runs:
            run.font.name = font_name
            run.font.size = Pt(font_size)
            run.font.color.rgb = color
        from pptx.util import Pt as _Pt
        p.space_after = _Pt(line_spacing_pt)
    return txBox


def navy_slide(slide):
    add_rect(slide, 0, 0, 13.33, 7.5, NAVY)


def add_footer(slide, text="Confidential — May 2026", slide_num=None):
    label = f"{text}   |   {slide_num}" if slide_num else text
    add_text(slide, label, 0.3, 7.1, 12, 0.35,
             font_size=9, color=MID_GRAY, align=PP_ALIGN.LEFT)


def accent_bar(slide, y=1.1, color=ORANGE):
    add_rect(slide, 0.5, y, 1.2, 0.06, color)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 1 — COVER
# ══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
navy_slide(slide)

# Decorative orange accent top-left block
add_rect(slide, 0, 0, 0.55, 7.5, ORANGE)

# Title
add_text(slide,
         "AI-Powered Trading Research\nfor the Individual Investor",
         0.9, 1.8, 9.5, 2.2,
         font_size=38, bold=True, color=WHITE, align=PP_ALIGN.LEFT)

# Tagline
add_text(slide,
         "Institutional-grade signal synthesis — finally accessible to individuals",
         0.9, 4.1, 9.5, 0.8,
         font_size=18, italic=True, color=ORANGE, align=PP_ALIGN.LEFT)

add_rect(slide, 0.9, 5.0, 5.5, 0.05, ORANGE)

add_text(slide, "Amit Garg  ·  amitgar@hotmail.com  ·  May 2026",
         0.9, 5.2, 9, 0.5,
         font_size=13, color=MID_GRAY, align=PP_ALIGN.LEFT)

add_text(slide, "CONFIDENTIAL",
         0.9, 5.8, 4, 0.4,
         font_size=11, bold=True, color=MID_GRAY, align=PP_ALIGN.LEFT)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 2 — THE PROBLEM
# ══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
navy_slide(slide)
accent_bar(slide, y=1.15)

add_text(slide, "The Problem", 0.5, 0.35, 12, 0.7,
         font_size=28, bold=True, color=WHITE)
add_text(slide, "Individual investors are making decisions with yesterday's tools",
         0.5, 1.25, 12, 0.5,
         font_size=16, italic=True, color=ORANGE)

# Three pain boxes
boxes = [
    ("Information\nAsymmetry",
     "Institutional desks have real-time signals:\nearnings drift, insider activity, options flow.\nRetail gets chart patterns."),
    ("Signal\nOverload",
     "1,000+ screeners. Thousands of technical\nindicators. No synthesis layer to say\nwhat actually matters today."),
    ("Execution\nGap",
     "Even when a retail trader identifies a setup,\nno automated discipline around entries,\nstops, or position sizing."),
]

for i, (title, body_text) in enumerate(boxes):
    x = 0.5 + i * 4.25
    add_rect(slide, x, 2.0, 3.9, 3.8, RGBColor(0x0F, 0x26, 0x50))
    add_rect(slide, x, 2.0, 3.9, 0.55, ORANGE)
    add_text(slide, title, x + 0.2, 2.05, 3.5, 0.5,
             font_size=15, bold=True, color=WHITE, align=PP_ALIGN.LEFT)
    add_text(slide, body_text, x + 0.2, 2.7, 3.5, 3.0,
             font_size=13, color=LIGHT_GRAY, align=PP_ALIGN.LEFT)

add_text(slide, "The result: 80% of individual traders underperform the index they're trying to beat.",
         0.5, 6.1, 12, 0.5,
         font_size=14, bold=True, color=ORANGE, align=PP_ALIGN.CENTER)

add_footer(slide, slide_num="2 / 12")


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 3 — MARKET OPPORTUNITY
# ══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
navy_slide(slide)
accent_bar(slide, y=1.15)

add_text(slide, "Market Opportunity", 0.5, 0.35, 12, 0.7,
         font_size=28, bold=True, color=WHITE)
add_text(slide, "A large, underserved market at the intersection of AI and retail investing",
         0.5, 1.25, 12, 0.5,
         font_size=16, italic=True, color=ORANGE)

# TAM / SAM / SOM circles (represented as boxes for reliability)
segments = [
    ("TAM",  "$47B",  "Global retail trading\nplatform market (2025)\n17.2% CAGR through 2030"),
    ("SAM",  "$8.2B", "Active algorithmic &\nquant retail traders\n(seeking signal + automation)"),
    ("SOM",  "$120M", "US-based active traders\nwilling to pay for\nAI research signals ($49–99/mo)"),
]

colors_seg = [RGBColor(0x0F, 0x26, 0x50), RGBColor(0x1A, 0x3A, 0x6A), ORANGE]
for i, (label, size, desc) in enumerate(segments):
    x = 0.6 + i * 4.1
    add_rect(slide, x, 2.0, 3.7, 4.0, colors_seg[i])
    add_text(slide, label, x + 0.2, 2.15, 3.3, 0.45,
             font_size=13, bold=True, color=MID_GRAY)
    add_text(slide, size, x + 0.2, 2.6, 3.3, 0.85,
             font_size=34, bold=True, color=WHITE)
    add_text(slide, desc, x + 0.2, 3.55, 3.3, 2.2,
             font_size=12, color=LIGHT_GRAY)

add_text(slide, "54M+ retail brokerage accounts in the US  ·  Robinhood alone: 24M+ active users",
         0.5, 6.25, 12.3, 0.5,
         font_size=13, bold=False, color=MID_GRAY, align=PP_ALIGN.CENTER)

add_footer(slide, slide_num="3 / 12")


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 4 — OUR SOLUTION
# ══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
navy_slide(slide)
accent_bar(slide, y=1.15)

add_text(slide, "Our Solution", 0.5, 0.35, 12, 0.7,
         font_size=28, bold=True, color=WHITE)
add_text(slide, "An AI research co-pilot that synthesizes real market signals into daily trade plans",
         0.5, 1.25, 12, 0.5,
         font_size=16, italic=True, color=ORANGE)

# Left: what it does
add_text(slide, "What It Does", 0.5, 2.0, 5.8, 0.45,
         font_size=16, bold=True, color=ORANGE)
lines_left = [
    "Monitors the market every morning before open",
    "Applies multiple independent signal filters",
    "Synthesizes signals with AI reasoning to rank setups",
    "Produces a daily briefing: top setups with entry, target, stop, and conviction",
    "Runs automated discipline: bracket orders, trailing exits, position limits",
    "Delivers an end-of-day performance audit with full decision trail",
]
add_multi_para(slide, lines_left, 0.5, 2.55, 5.8, 4.2,
               font_size=13, color=LIGHT_GRAY, bullet_char="✓")

# Right: what it is NOT
add_text(slide, "What It Is NOT", 7.0, 2.0, 5.8, 0.45,
         font_size=16, bold=True, color=MID_GRAY)
lines_right = [
    "Not a black box — every decision is auditable",
    "Not another chart screener",
    "Not AI generating random trade ideas",
    "Not a robo-advisor managing a passive portfolio",
    "Not dependent on a single signal or model",
]
add_multi_para(slide, lines_right, 7.0, 2.55, 5.8, 4.0,
               font_size=13, color=MID_GRAY, bullet_char="✗")

# Divider
add_rect(slide, 6.67, 2.0, 0.04, 4.8, RGBColor(0x33, 0x44, 0x66))

add_footer(slide, slide_num="4 / 12")


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 5 — HOW IT WORKS (no mechanics)
# ══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
navy_slide(slide)
accent_bar(slide, y=1.15)

add_text(slide, "How It Works", 0.5, 0.35, 12, 0.7,
         font_size=28, bold=True, color=WHITE)
add_text(slide, "A five-stage pipeline — from market open to end-of-day audit",
         0.5, 1.25, 12, 0.5,
         font_size=16, italic=True, color=ORANGE)

stages = [
    ("1", "Signal\nCapture",    "Pre-market: multiple independent market signals gathered automatically"),
    ("2", "AI\nSynthesis",      "AI reasoning layer weighs confluence of signals, filters low-conviction setups"),
    ("3", "Risk\nValidation",   "Position sizing, sector concentration, and capital guardrails enforced automatically"),
    ("4", "Execution",          "Bracketed entries with automated trailing exits — no emotional override"),
    ("5", "Audit &\nLearning",  "Full decision trail logged. Daily performance audit. Signal quality tracked over time."),
]

for i, (num, stage, desc) in enumerate(stages):
    x = 0.4 + i * 2.5
    add_rect(slide, x, 2.1, 2.25, 0.7, ORANGE)
    add_text(slide, num, x + 0.1, 2.12, 0.5, 0.65,
             font_size=22, bold=True, color=WHITE)
    add_text(slide, stage, x + 0.55, 2.12, 1.6, 0.65,
             font_size=13, bold=True, color=WHITE)
    add_rect(slide, x, 2.8, 2.25, 3.6, RGBColor(0x0F, 0x26, 0x50))
    add_text(slide, desc, x + 0.15, 2.95, 1.95, 3.3,
             font_size=12, color=LIGHT_GRAY)

    # Arrow between stages
    if i < 4:
        add_text(slide, "→", x + 2.25, 2.28, 0.25, 0.4,
                 font_size=18, bold=True, color=ORANGE, align=PP_ALIGN.CENTER)

add_text(slide, "Fully automated. Runs on a schedule. No human intervention required during market hours.",
         0.5, 6.55, 12.3, 0.5,
         font_size=13, bold=True, color=ORANGE, align=PP_ALIGN.CENTER)

add_footer(slide, slide_num="5 / 12")


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 6 — THE MOAT
# ══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
navy_slide(slide)
accent_bar(slide, y=1.15)

add_text(slide, "The MOAT", 0.5, 0.35, 12, 0.7,
         font_size=28, bold=True, color=WHITE)
add_text(slide, "Four compounding advantages — none individually unique, combination is",
         0.5, 1.25, 12, 0.5,
         font_size=16, italic=True, color=ORANGE)

moats = [
    ("Real Signals,\nNot Technicals",
     "Uses academically-documented market anomalies with known predictive power — not RSI crossovers that every retail screener already runs."),
    ("AI as\nSynthesizer",
     "AI doesn't generate signals — it synthesizes across multiple independent signals. When PEAD + institutional activity + momentum align, AI produces a conviction-ranked plan nobody else is generating."),
    ("Full\nAuditability",
     "Every decision is logged with the reasoning chain. Users can read exactly why each trade was selected. No black box — this is the foundation of trust and a future regulatory moat."),
    ("Operational\nRigor",
     "Health checks, guardrails, sector limits, exit discipline, performance audits. The system runs without human monitoring and fails safely. Competitors have none of this infrastructure."),
]

for i, (title, desc) in enumerate(moats):
    row = i // 2
    col = i % 2
    x = 0.5 + col * 6.4
    y = 2.0 + row * 2.5
    add_rect(slide, x, y, 6.1, 2.2, RGBColor(0x0F, 0x26, 0x50))
    add_rect(slide, x, y, 0.12, 2.2, ORANGE)
    add_text(slide, title, x + 0.3, y + 0.15, 2.5, 0.85,
             font_size=15, bold=True, color=WHITE)
    add_text(slide, desc, x + 0.3, y + 0.95, 5.6, 1.1,
             font_size=12, color=LIGHT_GRAY)

add_footer(slide, slide_num="6 / 12")


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 7 — COMPETITIVE LANDSCAPE
# ══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
navy_slide(slide)
accent_bar(slide, y=1.15)

add_text(slide, "Competitive Landscape", 0.5, 0.35, 12, 0.7,
         font_size=28, bold=True, color=WHITE)
add_text(slide, "Nobody owns the combination we've built",
         0.5, 1.25, 12, 0.5,
         font_size=16, italic=True, color=ORANGE)

# Table header
cols_x   = [0.5, 2.9, 5.05, 7.2, 9.35, 11.5]
col_w    = [2.3, 2.1, 2.1,  2.1,  2.1,   1.5]
headers  = ["", "Real Signals", "AI Synthesis", "Full Audit Trail", "Auto Execution", "Operational Infra"]
rows_data = [
    ("Composer / Streak",  "✗", "✗", "✗", "Partial", "✗"),
    ("QuantConnect",       "✗", "✗", "Partial", "✓", "✗"),
    ("TradingAgents",      "✗", "Partial", "✗", "✗", "✗"),
    ("Alpaca MCP",         "✗", "Natural lang", "✗", "✓", "✗"),
    ("This System",        "✓", "✓", "✓", "✓", "✓"),
]

header_y = 2.0
add_rect(slide, 0.4, header_y, 12.5, 0.5, RGBColor(0x0F, 0x26, 0x50))
for j, (hdr, cx, cw) in enumerate(zip(headers, cols_x, col_w)):
    add_text(slide, hdr, cx, header_y + 0.08, cw, 0.38,
             font_size=11, bold=True, color=ORANGE if j > 0 else WHITE,
             align=PP_ALIGN.CENTER if j > 0 else PP_ALIGN.LEFT)

for i, row in enumerate(rows_data):
    y = 2.55 + i * 0.72
    bg = ORANGE if row[0] == "This System" else RGBColor(0x12, 0x2A, 0x55) if i % 2 else RGBColor(0x0F, 0x24, 0x4A)
    add_rect(slide, 0.4, y, 12.5, 0.67, bg)
    for j, (val, cx, cw) in enumerate(zip(row, cols_x, col_w)):
        cell_color = WHITE if row[0] == "This System" else (
            GREEN if val == "✓" else RGBColor(0xC0, 0x39, 0x2B) if val == "✗" else LIGHT_GRAY
        )
        add_text(slide, val, cx, y + 0.1, cw, 0.5,
                 font_size=12, bold=(row[0] == "This System"),
                 color=cell_color,
                 align=PP_ALIGN.CENTER if j > 0 else PP_ALIGN.LEFT)

add_footer(slide, slide_num="7 / 12")


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 8 — TRACTION & VALIDATION
# ══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
navy_slide(slide)
accent_bar(slide, y=1.15)

add_text(slide, "Traction & Validation", 0.5, 0.35, 12, 0.7,
         font_size=28, bold=True, color=WHITE)
add_text(slide, "Paper validation underway — real capital deployment gate: June 2026",
         0.5, 1.25, 12, 0.5,
         font_size=16, italic=True, color=ORANGE)

# KPI boxes
kpis = [
    ("90%",  "Win-day rate\n(last 30 days)"),
    ("$716", "Avg daily P&L\non paper ($100K)"),
    ("0.78", "ML signal\nAUC score"),
    ("v5.7", "Production\nversion live"),
]

for i, (val, label) in enumerate(kpis):
    x = 0.5 + i * 3.1
    add_rect(slide, x, 2.0, 2.8, 2.0, RGBColor(0x0F, 0x26, 0x50))
    add_rect(slide, x, 2.0, 2.8, 0.06, ORANGE)
    add_text(slide, val, x + 0.15, 2.2, 2.5, 1.0,
             font_size=40, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_text(slide, label, x + 0.15, 3.2, 2.5, 0.7,
             font_size=12, color=MID_GRAY, align=PP_ALIGN.CENTER)

# Milestones
add_text(slide, "System Milestones", 0.5, 4.2, 12, 0.45,
         font_size=15, bold=True, color=ORANGE)

milestones = [
    ("✓ LIVE", "Fully automated pipeline — premarket, intraday, end-of-day — running on GitHub Actions schedule"),
    ("✓ LIVE", "Paper trading on Alpaca broker API — bracket orders with native trailing stop loss"),
    ("✓ LIVE", "Real-time dashboard: signal quality, position monitoring, agent performance scorecard"),
    ("✓ LIVE", "ML signal ranker trained and integrated — AUC 0.78 on held-out data"),
    ("⏳ GATE", "June 1, 2026: 2-week live validation gate — pass criteria before real capital deployment"),
]

for i, (status, desc) in enumerate(milestones):
    y = 4.75 + i * 0.36
    color = GREEN if "LIVE" in status else ORANGE
    add_text(slide, status, 0.5, y, 1.3, 0.33,
             font_size=11, bold=True, color=color)
    add_text(slide, desc, 1.9, y, 10.9, 0.33,
             font_size=11, color=LIGHT_GRAY)

add_footer(slide, slide_num="8 / 12")


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 9 — BUSINESS MODEL
# ══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
navy_slide(slide)
accent_bar(slide, y=1.15)

add_text(slide, "Business Model", 0.5, 0.35, 12, 0.7,
         font_size=28, bold=True, color=WHITE)
add_text(slide, "Three paths — sequenced by validation confidence",
         0.5, 1.25, 12, 0.5,
         font_size=16, italic=True, color=ORANGE)

paths = [
    ("Path A",
     "Personal Alpha",
     "Now → 2027",
     ORANGE,
     [
         "$100K paper → real capital once June gate passes",
         "Goal: $500–700/day consistent cash flow",
         "55% win rate × 3:1 R:R = $48 EV per trade",
         "15 trades/day = $720 expected daily value",
         "Scale via compounding — same % on growing capital",
     ]),
    ("Path B",
     "LP Fund Structure",
     "2027 → 2028",
     RGBColor(0x1A, 0x7A, 0x3A),
     [
         "12 months audited live returns required first",
         "Raise $500K–$2M from accredited investors",
         "Below $150M AUM — ERA exemption applies",
         "20% carry on profits",
         "Sharpe ≥ 1.5 over 200+ live trades as entry gate",
     ]),
    ("Path C",
     "SaaS Signal Product",
     "2027 → beyond",
     RGBColor(0x6A, 0x1A, 0x6A),
     [
         "Daily AI research digest — top 3–5 setups with reasoning",
         "$49–$99/month per subscriber",
         "1,000 subscribers = $600K–$1.2M ARR",
         "No broker-dealer registration needed (research product)",
         "Leverages existing signal engine + dashboard",
     ]),
]

for i, (label, title, timeline, color, bullets) in enumerate(paths):
    x = 0.4 + i * 4.3
    add_rect(slide, x, 2.0, 4.0, 4.9, RGBColor(0x0F, 0x26, 0x50))
    add_rect(slide, x, 2.0, 4.0, 0.55, color)
    add_text(slide, label, x + 0.2, 2.05, 1.2, 0.45,
             font_size=11, bold=True, color=WHITE)
    add_text(slide, timeline, x + 1.5, 2.1, 2.3, 0.38,
             font_size=10, color=WHITE, align=PP_ALIGN.RIGHT)
    add_text(slide, title, x + 0.2, 2.65, 3.6, 0.5,
             font_size=16, bold=True, color=WHITE)
    for j, b in enumerate(bullets):
        add_text(slide, f"• {b}", x + 0.2, 3.25 + j * 0.68, 3.6, 0.62,
                 font_size=11, color=LIGHT_GRAY)

add_footer(slide, slide_num="9 / 12")


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 10 — ROAD MAP
# ══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
navy_slide(slide)
accent_bar(slide, y=1.15)

add_text(slide, "Road Map", 0.5, 0.35, 12, 0.7,
         font_size=28, bold=True, color=WHITE)
add_text(slide, "Phased — each phase gate-locked to validated outcomes",
         0.5, 1.25, 12, 0.5,
         font_size=16, italic=True, color=ORANGE)

phases = [
    ("Phase 0",  "Now → Jun 1",  "Validate execution",
     ["Native trailing stop confirmed live", "Win rate ≥ 80% on paper", "VWAP signal quality measured", "No integrity failures"]),
    ("Phase 1",  "Jun → Sep 2026", "Expand signal layer",
     ["Post-earnings drift scanner", "Insider buying feed (SEC EDGAR)", "Unusual options activity", "Multi-day hold mode"]),
    ("Phase 2",  "Sep → Dec 2026", "Real capital",
     ["Deploy $10K–$25K real money", "Track paper-to-live degradation", "Validate live fill quality", "Scale if win rate holds ±10%"]),
    ("Phase 3",  "2027+",  "Product or fund",
     ["Choose Path A / B / C based on live Sharpe", "If SaaS: build subscription UI", "If fund: legal structure + LP reporting", "If personal: compound capital"]),
]

phase_colors = [ORANGE, RGBColor(0x1A, 0x7A, 0x3A), RGBColor(0x2A, 0x5A, 0xA0), RGBColor(0x6A, 0x1A, 0x6A)]

for i, (phase, timeline, goal, items) in enumerate(phases):
    x = 0.4 + i * 3.2
    add_rect(slide, x, 2.0, 3.0, 0.65, phase_colors[i])
    add_text(slide, phase, x + 0.15, 2.05, 2.7, 0.35,
             font_size=14, bold=True, color=WHITE)
    add_text(slide, timeline, x + 0.15, 2.38, 2.7, 0.28,
             font_size=10, color=WHITE)
    add_rect(slide, x, 2.65, 3.0, 0.45, RGBColor(0x0F, 0x26, 0x50))
    add_text(slide, goal, x + 0.15, 2.7, 2.7, 0.38,
             font_size=12, bold=True, color=ORANGE)
    add_rect(slide, x, 3.1, 3.0, 3.7, RGBColor(0x0A, 0x1E, 0x40))
    for j, item in enumerate(items):
        add_text(slide, f"• {item}", x + 0.15, 3.2 + j * 0.82, 2.7, 0.75,
                 font_size=11, color=LIGHT_GRAY)

    if i < 3:
        add_text(slide, "→", x + 3.0, 2.9, 0.2, 0.4,
                 font_size=18, bold=True, color=MID_GRAY, align=PP_ALIGN.CENTER)

add_footer(slide, slide_num="10 / 12")


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 11 — FOUNDER
# ══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
navy_slide(slide)
accent_bar(slide, y=1.15)

add_text(slide, "The Builder", 0.5, 0.35, 12, 0.7,
         font_size=28, bold=True, color=WHITE)
add_text(slide, "Product + engineering + operational depth — rare combination for this problem",
         0.5, 1.25, 12, 0.5,
         font_size=16, italic=True, color=ORANGE)

# Left column — bio
add_text(slide, "Amit Garg", 0.5, 2.0, 6, 0.6,
         font_size=22, bold=True, color=WHITE)
add_text(slide, "VP / Head of Product  ·  25 years in tech",
         0.5, 2.6, 6, 0.4,
         font_size=14, color=MID_GRAY)

bio_items = [
    "ServiceNow 2012–2026: built $100M+ product 0-to-1 in 12 months; scaled team 0→60+",
    "Owned internal control plane: 45+ data centers, 99.99% availability, 1.8K enterprise customers",
    "Microsoft 1999–2012; B.Tech CS — started as a developer, reads and writes code",
    "AIOps, observability, data platforms, automated workflows — exactly this problem domain",
    "Certified AI Product Manager  ·  PMP  ·  Two ServiceNow Technology Awards for Innovation",
]
for i, item in enumerate(bio_items):
    add_text(slide, f"• {item}", 0.5, 3.15 + i * 0.62, 6.1, 0.58,
             font_size=12, color=LIGHT_GRAY)

# Right column — why uniquely qualified
add_rect(slide, 7.0, 2.0, 5.9, 4.9, RGBColor(0x0F, 0x26, 0x50))
add_rect(slide, 7.0, 2.0, 5.9, 0.06, ORANGE)

add_text(slide, "Why Uniquely Qualified", 7.2, 2.15, 5.5, 0.5,
         font_size=15, bold=True, color=ORANGE)

quals = [
    ("Domain",     "Built production-grade automated systems at scale — not a side project"),
    ("Product",    "Knows how to take a research prototype to a product users pay for"),
    ("Technical",  "Wrote the agent pipeline, ML model, and broker integration end-to-end"),
    ("Investing",  "Personal capital at stake — aligned incentive structure"),
    ("Network",    "Enterprise product and SRE network — natural early adopter base for Scenario C"),
]
for i, (cat, desc) in enumerate(quals):
    y = 2.75 + i * 0.75
    add_text(slide, cat, 7.2, y, 1.5, 0.5,
             font_size=12, bold=True, color=ORANGE)
    add_text(slide, desc, 8.85, y, 3.9, 0.65,
             font_size=12, color=LIGHT_GRAY)

add_footer(slide, slide_num="11 / 12")


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 12 — THE ASK
# ══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
navy_slide(slide)

# Full-width orange top band
add_rect(slide, 0, 0, 13.33, 1.0, ORANGE)
add_text(slide, "The Ask", 0.5, 0.18, 12, 0.65,
         font_size=30, bold=True, color=WHITE)

add_text(slide, "We're at an inflection point. The infrastructure is built. The validation gate opens June 1.",
         0.5, 1.3, 12.3, 0.55,
         font_size=16, color=LIGHT_GRAY, align=PP_ALIGN.CENTER)

# Three ask boxes
asks = [
    ("Strategic\nConversation",
     "If you've seen similar systems scale — what breaks between paper and real money? We want to talk to people who've navigated this."),
    ("Early Access\n& Feedback",
     "Looking for 5–10 active traders to use the daily briefing and provide signal quality feedback before we open more broadly."),
    ("Investment\nInterest",
     "If traction metrics hold after June gate: raising a seed round to accelerate signal layer expansion and SaaS product build."),
]

for i, (title, body_text) in enumerate(asks):
    x = 0.5 + i * 4.2
    add_rect(slide, x, 2.05, 3.9, 3.7, RGBColor(0x0F, 0x26, 0x50))
    add_rect(slide, x, 2.05, 3.9, 0.55, NAVY)
    add_rect(slide, x, 2.05, 0.1, 3.7, ORANGE)
    add_text(slide, title, x + 0.25, 2.1, 3.5, 0.5,
             font_size=15, bold=True, color=ORANGE)
    add_text(slide, body_text, x + 0.25, 2.72, 3.55, 2.8,
             font_size=13, color=LIGHT_GRAY)

# Contact
add_rect(slide, 0.5, 5.95, 12.3, 0.06, ORANGE)
add_text(slide, "Amit Garg   ·   amitgar@hotmail.com   ·   linkedin.com/in/mramitgarg   ·   Sammamish, WA",
         0.5, 6.1, 12.3, 0.5,
         font_size=14, color=WHITE, align=PP_ALIGN.CENTER)

add_text(slide, "Confidential — not for distribution  ·  May 2026",
         0.5, 6.7, 12.3, 0.4,
         font_size=10, color=MID_GRAY, align=PP_ALIGN.CENTER)


# ── Save ──────────────────────────────────────────────────────────────────────
_project_dir = "/Users/amitgarg/Claude Projects/trading-agent"
out = os.path.join(_project_dir, "Trading_Agent_Pitch_Deck.pptx")
prs.save(out)
print(f"Saved: {out}")
